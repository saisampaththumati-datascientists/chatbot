import streamlit as st
from PyPDF2 import PdfReader
from llama_index.llms.ollama import Ollama
from youtube_transcript_api import YouTubeTranscriptApi, TranscriptsDisabled
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import requests
from bs4 import BeautifulSoup
import numpy as np
import re 
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from typing import List
# Initialize Ollama LLM
ollama_llm = Ollama(model="llama3.2")

# Function to check if the URL is a YouTube link
def is_youtube_url(url):
    return "youtube.com" in url or "youtu.be" in url

# Function to extract text from a YouTube video
def extract_text_from_youtube(url):
    try:
        if "youtube.com" in url:
            video_id = url.split("v=")[1].split("&")[0]
        elif "youtu.be" in url:
            video_id = url.split("/")[-1]

        transcript = YouTubeTranscriptApi.get_transcript(video_id)
        text = " ".join([item['text'] for item in transcript])
        return text, 'en'
    except TranscriptsDisabled:
        return "Transcripts are disabled for this video.", None
    except Exception as e:
        if 'transcript' in str(e):
            available_transcripts = YouTubeTranscriptApi.list_transcripts(video_id)
            available_languages = [t.language for t in available_transcripts]
            return f"No transcripts found in English. Available languages: {available_languages}", None
        return f"Error extracting YouTube transcript: {str(e)}", None

# Function to extract text from a PDF file
def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""  # Handle None returned by extract_text
    return text

# Function to extract text from a generic URL
def extract_text_from_url(url):
    try:
        response = requests.get(url)
        response.raise_for_status()  # Raise an error for bad responses
        soup = BeautifulSoup(response.text, 'html.parser')

        # Extract text from paragraph tags
        paragraphs = soup.find_all('p')
        text = " ".join([para.get_text() for para in paragraphs])
        return text, 'en'  # Default to English
    except Exception as e:
        return f"Error extracting text from URL: {str(e)}", None

# Function to process text with LLM
def process_chunk_with_llm(chunk, task_type="summary"):
    try:
        prompt_map = {
            "summary": f"Summarize this text: {chunk}",
            "translate": f"Translate this text to English: {chunk}",
            "heading": f"Create a concise single heading for this content: {chunk}"
        }
        response = ollama_llm.complete(prompt=prompt_map.get(task_type, "Unsupported task type."))
        
        if isinstance(response, dict) and "text" in response:
            return response["text"]
        elif isinstance(response, str):
            return response
        else:
            return f"Unexpected response format: {response}"
    except Exception as e:
        return f"Error processing chunk with LLM: {str(e)}"

# Function to split text into chunks for processing
def split_text_into_chunks(text, chunk_size=500):
    words = text.split()
    chunks = [' '.join(words[i:i + chunk_size]) for i in range(0, len(words), chunk_size)]
    return chunks

# Function to create a downloadable .txt file
def create_downloadable_file(content, filename="output.txt"):
    with open(filename, "w") as file:
        file.write(content)
    return filename


def split_text(text, max_chars=500):
    """
    Splits the text into smaller chunks of a specified maximum number of characters.

    Args:
        text (str): The text to be split.
        max_chars (int): Maximum number of characters in each chunk.

    Returns:
        List[str]: List of text chunks.
    """
    words = text.split(" ")
    chunks = []
    chunk = ""

    for word in words:
        # Check if adding the next word would exceed the maximum character count
        if len(chunk) + len(word) + 1 <= max_chars:
            chunk += word + " "
        else:
            if chunk:  # Ensure the chunk is not empty before adding
                chunks.append(chunk.strip())  # Add the chunk and reset
            chunk = word + " "

    if chunk:  # Add any remaining text as the last chunk
        chunks.append(chunk.strip())

    return chunks

def create_ppt_from_text(content, headings, filename="presentation.pptx", heading_font_size=20, body_font_size=14):
    """
    Creates a PowerPoint presentation with formatted headings and text content.

    Args:
        content (List[str]): List of text content to be added to the slides.
        headings (List[str]): List of headings for each slide.
        filename (str): Name of the file to save the presentation as.
        heading_font_size (int): Font size for slide headings.
        body_font_size (int): Font size for body content.

    Returns:
        str: The filename of the saved PowerPoint presentation.
    """
    prs = Presentation()

    # Iterate through each heading and corresponding text content
    for heading, text_chunk in zip(headings, content):
        # Split each text chunk into smaller sub-chunks to avoid overcrowding
        sub_chunks = split_text(text_chunk, max_chars=500)

        # Create slides for each sub-chunk
        for sub_chunk in sub_chunks:
            if sub_chunk:  # Ensure the sub_chunk is not empty
                # Add a new slide with a blank layout
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use blank layout

                # Create and format the title (heading)
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1.0))
                title_frame = title_box.text_frame
                title_frame.text = heading
                title_format = title_frame.paragraphs[0].font
                title_format.size = Pt(heading_font_size)
                title_format.bold = True
                title_format.color.rgb = RGBColor(0, 51, 102)  # Dark blue color for headings

                # Create and format the content (body text)
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.5), Inches(5.0))
                content_frame = content_box.text_frame
                content_frame.text = sub_chunk
                content_format = content_frame.paragraphs[0].font
                content_format.size = Pt(body_font_size)
                content_format.color.rgb = RGBColor(0, 0, 0)  # Black color for body text

                # Adjust text alignment and wrapping if needed
                for paragraph in content_frame.paragraphs:
                    paragraph.word_wrap = True  # Enable word wrapping to avoid text overflow
                    paragraph.alignment = None  # Set alignment to default (left)

    # Save the presentation to a file
    prs.save(filename)
    return filename



def find_best_matching_chunk(question, chunks):
    vectorizer = TfidfVectorizer()
    vectors = vectorizer.fit_transform([question] + chunks)
    cosine_similarities = cosine_similarity(vectors[0:1], vectors[1:]).flatten()
    best_match_index = np.argmax(cosine_similarities)
    return chunks[best_match_index]

# Define the layout and options for the first page
def page1():
    st.title("LLM Text Processor")

    # Initialize session state variables if they don't exist
    if "process_disabled" not in st.session_state:
        st.session_state.process_disabled = True
    if "processed_results" not in st.session_state:
        st.session_state.processed_results = []
    if "chunks" not in st.session_state:
        st.session_state.chunks = []
    if "headings" not in st.session_state:
        st.session_state.headings = []

    st.header("User Input Options")

    # Provide unique keys for each input element
    url_input = st.text_input("Enter the URL for a web page:", key="url_input")  # General URL input
    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf", key="pdf_file")
    option = st.radio("Select what you want to do:", ('Summary'), key="option_radio")  # Added more options
    file_type = st.selectbox("Select the output file format:", ["Text (.txt)", "PowerPoint (.pptx)"], key="output_format")

    # Enable process button only if a URL or file is provided
    st.session_state.process_disabled = not (url_input or uploaded_file)

    process_button = st.button("Process", disabled=st.session_state.process_disabled)

    if process_button and not st.session_state.process_disabled:
        st.session_state.process_disabled = True
        text_data = ""

        if url_input:
            with st.spinner("Extracting text from the URL..."):
                if is_youtube_url(url_input):
                    text_data, _ = extract_text_from_youtube(url_input)
                else:
                    text_data, _ = extract_text_from_url(url_input)

        if uploaded_file:
            with st.spinner("Extracting text from PDF file..."):
                text_data = extract_text_from_pdf(uploaded_file)

        if text_data:
            st.write(f"Extracted Text Data Length: {len(text_data)} characters")
            chunks = split_text_into_chunks(text_data, chunk_size=500)
            st.write(f"Number of chunks created: {len(chunks)}")

            processed_results = []
            headings = []

            for chunk in chunks:
                with st.spinner(f"Processing chunk for {option}..."):
                    result = process_chunk_with_llm(chunk, task_type=option.lower())
                    cleaned_result = re.sub(r"Unexpected response format: ", "", result)
                    if isinstance(cleaned_result, str):
                        processed_results.append(cleaned_result)
                        # Only generate headings if the task is not to generate headings itself
                        if option.lower() != "generate headings":
                            heading = process_chunk_with_llm(chunk, task_type="heading")
                            cleaned_heading = re.sub(r"Unexpected response format: ", "", heading)
                            headings.append(cleaned_heading if isinstance(cleaned_heading, str) else "Untitled Section")

            st.session_state.processed_results = processed_results
            st.session_state.chunks = chunks
            st.session_state.headings = headings

            # Display processed results
            if processed_results:
                st.write("## Processed Output")
                for heading, result in zip(headings, processed_results):
                    st.write(f"**{heading}:** {result}")

                # Create a downloadable file with all processed chunks
                question_answer_filename = create_downloadable_file("\n\n".join(st.session_state.chunks), filename="chunked_data.txt")
                with open(question_answer_filename, "rb") as file:
                    st.download_button(
                        label="Download Chunked Data File",
                        data=file,
                        file_name=question_answer_filename,
                        mime="text/plain"
                    )

            # Create downloadable file based on selected format
            if file_type == "Text (.txt)":
                filename = create_downloadable_file("\n\n".join(processed_results), filename=f"processed_{option.lower()}.txt")
            else:
                filename = create_ppt_from_text("\n\n".join(processed_results), st.session_state.headings, filename="processed_presentation.pptx")

            with open(filename, "rb") as file:
                st.download_button(
                    label=f"Download {file_type.split()[0]} File",
                    data=file,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation" if file_type == "PowerPoint (.pptx)" else "text/plain"
                )

        st.session_state.process_disabled = False
